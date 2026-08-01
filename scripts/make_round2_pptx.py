"""Generate the Round-2 business deck (docs/PRESENTATION_ROUND2.pptx).

Round 1 was "does it work?". The panel answered that themselves - they liked
the idea and the build, then stopped us on time and asked seven business
questions. This deck answers only those seven.

Structured as FOUR parts, each with a divider slide, so a four-way handover is
obvious to the panel without ever printing a speaker's name on a slide:

    Part 1  Business strategy  - who will buy it
    Part 2  Technology         - camera, AI, data, maps, detection flow
    Part 3  Money              - pricing, subscription, revenue
    Part 4  Implementation     - rollout, risks, the ask

Part 2 is the longest on purpose. "Where does the zone data come from?" is the
question the panel pressed hardest on, and answering it properly means showing
the whole path from a camera on a pole to a challan in an officer's queue.

Same design rules as the main deck:
  * Light theme, one teal accent - readable over a compressed Meet stream.
  * Nothing backstage on a slide. No speaker names, no timings.
  * Every commercial figure is labelled as a modelled assumption, with the
    working shown. A number we cannot source does not go on a slide.

    python scripts/make_round2_pptx.py
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from _pptx_helpers import (ACCENT, AMBER, GREEN, H, INK, MUTED, PANEL,
                           PANEL_LINE, W, accent_bar, bg, box, card, faded_bg,
                           header, stat)

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
os.chdir(ROOT)

PAPER = RGBColor(0xFF, 0xFF, 0xFF)


def _rect(slide, x, y, w, h, rgb, line=None):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = rgb
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def table(s, data, widths, top=2.1, left=0.85, total=11.65, rh=0.5, size=11.5):
    """Row-banded table. data[0] is the header row.

    python-pptx's native table styling is fought more easily than configured,
    so this draws rectangles and textboxes - which also keeps the look
    identical to the cards used everywhere else in the deck.
    """
    y = Inches(top)
    for r, row in enumerate(data):
        fill = ACCENT if r == 0 else (PANEL if r % 2 == 0 else PAPER)
        _rect(s, Inches(left), y, Inches(total), Inches(rh), fill,
              line=None if r == 0 else PANEL_LINE)
        x = left
        for c, cell in enumerate(row):
            cw = total * widths[c]
            box(s, Inches(x + 0.14), y + Inches(0.09), Inches(cw - 0.24),
                Inches(rh), str(cell), size=size,
                colour=PAPER if r == 0 else INK, bold=(r == 0), space=0)
            x += cw
        y += Inches(rh)
    return y


def note(s, text, y=6.85, colour=ACCENT, size=14, bold=True):
    box(s, Inches(0.85), Inches(y), Inches(11.65), Inches(0.6), text,
        size=size, colour=colour, bold=bold)


def rows(s, items, top=2.6, gap=0.85, hgt=0.72, label_w=3.9, size=13.5):
    """Label-on-the-left, explanation-on-the-right list of cards."""
    for i, (label, detail) in enumerate(items):
        y = Inches(top + i * gap)
        card(s, Inches(0.85), y, Inches(11.65), Inches(hgt))
        box(s, Inches(1.1), y + Inches(0.13), Inches(label_w), Inches(0.45),
            label, size=size, colour=ACCENT, bold=True)
        box(s, Inches(1.2 + label_w), y + Inches(0.15), Inches(11.1 - label_w),
            Inches(0.45), detail, size=12, colour=INK, space=0)


def divider(s, number, title, answers):
    """Part divider. Deliberately plain - it is a breath, not a slide."""
    bg(s)
    accent_bar(s)
    _rect(s, Inches(0.85), Inches(2.55), Inches(0.1), Inches(2.0), ACCENT)
    box(s, Inches(1.25), Inches(2.5), Inches(10.5), Inches(0.5),
        f"PART {number}", size=15, colour=ACCENT, bold=True)
    box(s, Inches(1.25), Inches(3.0), Inches(11.0), Inches(1.0), title,
        size=40, colour=INK, bold=True)
    box(s, Inches(1.25), Inches(4.05), Inches(11.0), Inches(0.5), answers,
        size=16, colour=MUTED, space=0)


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]
    S = lambda: prs.slides.add_slide(blank)   # noqa: E731

    # ============================================================ 1  title
    s = S(); bg(s)
    faded_bg(s, "12_junction_cctv_clean.jpg", opacity=0.15)
    accent_bar(s)
    box(s, Inches(0.9), Inches(1.75), Inches(11.5), Inches(0.6),
        "SECOND ROUND", size=15, colour=ACCENT, bold=True)
    box(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.7),
        "From a system that works to a business that ships",
        size=42, colour=INK, bold=True)
    box(s, Inches(0.9), Inches(4.05), Inches(11.5), Inches(0.7),
        "AI Smart Traffic Intelligence Platform  -  how it gets onto Sri Lanka's roads",
        size=18, colour=ACCENT)
    box(s, Inches(0.9), Inches(5.5), Inches(11.5), Inches(0.5),
        "Team Three Hacks", size=17, colour=INK, bold=True)
    box(s, Inches(0.9), Inches(6.0), Inches(11.5), Inches(0.5),
        "Startup Innovation Competition 2026", size=14, colour=MUTED)

    # ============================================================ 2  the brief
    s = S(); bg(s)
    header(s, "You asked us seven questions", "WHAT THIS DECK ANSWERS")
    table(s, [
        ["The question you put to us", "Part", "Slide"],
        ["Who is the customer, and how do you sell it?", "1", "4"],
        ["How do you connect with the Sri Lanka Police?", "1", "5-7"],
        ["How do you import data from existing CCTV?", "2", "10-11"],
        ["At scale, where does the zone data come from?", "2", "13-15"],
        ["Revenue model - your cost, your earnings from one province", "3", "18-23"],
        ["SWOT analysis", "4", "27"],
        ["How do you implement it in the real world?", "4", "25"],
        ["What challenges will you face?", "4", "26"],
    ], widths=[0.74, 0.13, 0.13], top=2.05, rh=0.5)
    note(s, "We have not re-pitched the detection. You already said it works. "
            "This is how it reaches a road.", y=6.4)

    # ############################################################ PART 1
    divider(S(), 1, "Business strategy",
            "Who will buy this, and how do we get to them?")

    # ------------------------------------------------------------ 4  customers
    s = S(); bg(s)
    header(s, "Who benefits is not who signs", "THE CUSTOMER MAP")
    table(s, [
        ["Segment", "Who signs", "Why they buy", "Cycle"],
        ["Municipal councils", "Municipal Commissioner",
         "Own CCTV, get no value from it", "1-3 months"],
        ["Sri Lanka Police - Traffic", "Ministry of Public Security",
         "Manpower gap; needs evidence", "12-24 months"],
        ["RDA / Expressways", "Expressway Division",
         "Speed enforcement already budgeted", "6-12 months"],
        ["Private fleets", "Fleet owner",
         "Insurance, liability, driver behaviour", "2-6 weeks"],
        ["Insurers", "Head of Motor", "Better risk pricing", "3-6 months"],
        ["Campuses, ports, BOI zones", "Facilities head",
         "Private roads - no legal complexity", "2-6 weeks"],
        ["Donor programmes", "Programme manager",
         "Road-safety targets with money attached", "6-12 months"],
    ], widths=[0.24, 0.22, 0.38, 0.16], top=2.05, rh=0.52, size=11)
    note(s, "The police are the biggest customer and the slowest customer. Three students "
            "cannot live two years on no revenue - so we run two tracks.", y=6.5, size=14)

    # ------------------------------------------------------------ 5  the ladder
    s = S(); bg(s)
    header(s, "We do not ask for a contract. We climb a ladder.",
           "GETTING TO THE SRI LANKA POLICE")
    ladder = [
        ("MONTH 0-3", "Give away evidence, not enforcement",
         "Free 30-day road-safety audit at one junction for one council. We issue ZERO fines. "
         "They get data they have never had. Nobody has to be brave to say yes."),
        ("MONTH 3-6", "Get a named champion",
         "Take that report to the National Council for Road Safety and Police Traffic HQ. "
         "Ask for one thing that costs them nothing: a letter of support."),
        ("MONTH 6-12", "Supervised enforcement pilot",
         "One division. Cameras we control. Every challan approved by an officer. "
         "We publish officer-time saved and the fall in the violation rate."),
        ("MONTH 12-24", "Procurement - through a partner",
         "Sub-contract under an integrator who already holds government vendor registration. "
         "Turns a two-year procurement problem into a two-month partnership."),
    ]
    y = 1.92
    for when, title, detail in ladder:
        card(s, Inches(0.85), Inches(y), Inches(11.65), Inches(1.14))
        box(s, Inches(1.05), Inches(y + 0.13), Inches(1.7), Inches(0.4), when,
            size=11, colour=ACCENT, bold=True)
        box(s, Inches(2.85), Inches(y + 0.09), Inches(9.4), Inches(0.4), title,
            size=15.5, colour=INK, bold=True)
        box(s, Inches(2.85), Inches(y + 0.52), Inches(9.4), Inches(0.55), detail,
            size=11.5, colour=MUTED, space=0)
        y += 1.24
    note(s, "Each rung is something the other side can say yes to without spending money "
            "or taking a risk.", y=6.88, size=13.5)

    # ------------------------------------------------------------ 6  the rule
    s = S(); bg(s)
    accent_bar(s)
    box(s, Inches(0.9), Inches(2.2), Inches(11.5), Inches(0.5),
        "THE DESIGN DECISION THAT MAKES ENFORCEMENT POSSIBLE", size=14,
        colour=ACCENT, bold=True)
    box(s, Inches(0.9), Inches(2.85), Inches(11.5), Inches(1.0),
        "The AI proposes. A police officer decides.", size=38, colour=INK, bold=True)
    box(s, Inches(0.9), Inches(4.1), Inches(11.5), Inches(0.9),
        "Every challan appears in an officer's review queue with the photograph, "
        "the rule and the confidence. It becomes a fine only when a human presses approve.",
        size=17, colour=INK, space=0)
    for i, (t, d) in enumerate([
            ("Legal", "Evidence Act No. 14 of 1995 - a human is accountable"),
            ("Trust", "No citizen is fined by a machine"),
            ("Accuracy", "Officer rejections measure our false-positive rate"),
            ("Politics", "Kills the 'a robot fined me' headline")]):
        x = Inches(0.85 + i * 2.95)
        card(s, x, Inches(5.55), Inches(2.75), Inches(1.15))
        box(s, x + Inches(0.18), Inches(5.67), Inches(2.4), Inches(0.4), t,
            size=14, colour=ACCENT, bold=True)
        box(s, x + Inches(0.18), Inches(6.05), Inches(2.4), Inches(0.7), d,
            size=10.5, colour=MUTED, space=0)

    # ------------------------------------------------------------ 7  two tracks
    s = S(); bg(s)
    header(s, "The private track funds the government track",
           "HOW WE STAY ALIVE FOR TWO YEARS")
    for i, (t, colour, items) in enumerate([
        ("GOVERNMENT TRACK  -  the prize", ACCENT, [
            "Councils, Police, RDA, expressways",
            "12-24 month cycle, tenders, bonds, audited accounts",
            "Buys legitimacy, scale and the reference every other market wants",
            "**Pays us in year 2, not year 1**"]),
        ("PRIVATE TRACK  -  the runway", GREEN, [
            "Fleets, universities, ports, factories, BOI zones",
            "Private roads: no procurement, no legal complexity",
            "A bus company can sign in a week",
            "**Pays us in month 2, and produces the accuracy record**"])]):
        x = Inches(0.85 + i * 6.0)
        card(s, x, Inches(2.1), Inches(5.65), Inches(3.5))
        box(s, x + Inches(0.28), Inches(2.32), Inches(5.1), Inches(0.5), t,
            size=15, colour=colour, bold=True)
        box(s, x + Inches(0.28), Inches(2.95), Inches(5.1), Inches(2.5), items,
            size=13, colour=INK, space=11)
    card(s, Inches(0.85), Inches(5.85), Inches(11.65), Inches(1.15))
    box(s, Inches(1.15), Inches(6.0), Inches(11.1), Inches(0.9),
        "**Highest-leverage channel:** get one motor insurer to discount premiums for "
        "fleets running our driver scorecard. Their agents then sell us to every fleet "
        "in their book, and we pay nothing for it.", size=14, colour=INK, space=0)

    # ############################################################ PART 2
    divider(S(), 2, "Technology",
            "Camera, AI, data, maps - how a moving vehicle becomes a challan")

    # ------------------------------------------------------------ 9  the flow
    s = S(); bg(s)
    header(s, "From a camera on a pole to a challan", "THE WHOLE SYSTEM IN SIX STEPS")
    steps = [
        ("1  CAMERA", "Existing council CCTV, our own camera, or a drone. "
                      "We pull the video over the ONVIF / RTSP standard."),
        ("2  EDGE BOX", "One small computer in the cabinet that is already there. "
                        "It handles about four cameras."),
        ("3  DETECT", "The AI finds vehicles, riders, helmets, phones and number plates, "
                      "and gives every vehicle an ID it keeps as it moves."),
        ("4  MEASURE", "The map calibration turns pixels into metres, so we get real "
                       "speed and real position on the road - not a guess."),
        ("5  JUDGE", "Ten rules run at once. Each one needs several frames of agreement "
                     "before it will accuse anybody."),
        ("6  EVIDENCE", "A photograph, stamped with the plate, time, place and rule, "
                        "goes to an officer's queue. Approved, it becomes a challan."),
    ]
    for i, (t, d) in enumerate(steps):
        x = Inches(0.85 + (i % 3) * 3.92)
        y = Inches(2.05 + (i // 3) * 1.85)
        card(s, x, y, Inches(3.72), Inches(1.65))
        box(s, x + Inches(0.22), y + Inches(0.14), Inches(3.3), Inches(0.4), t,
            size=14, colour=ACCENT, bold=True)
        box(s, x + Inches(0.22), y + Inches(0.55), Inches(3.32), Inches(1.0), d,
            size=11, colour=INK, space=0)
    note(s, "Only step 6 ever leaves the site. Everything before it happens on the pole.",
         y=5.9, size=16)

    # ------------------------------------------------------------ 10  cctv in
    s = S(); bg(s)
    header(s, "Cameras that are already on the poles", "STEP 1  -  WHERE THE PICTURES COME FROM")
    box(s, Inches(0.85), Inches(1.98), Inches(11.65), Inches(0.6),
        "Almost every camera and recorder installed in this country speaks **ONVIF over "
        "RTSP**. We pull a stream. No firmware change, no rewiring, nothing new on the pole.",
        size=14.5, colour=INK, space=0)
    for i, (t, d) in enumerate([
            ("Council CCTV", "Already installed and paid for. We just read it."),
            ("Our own camera", "Where none exists, or the angle is wrong."),
            ("Drone", "Already working today - we fly one.")]):
        x = Inches(0.85 + i * 3.92)
        card(s, x, Inches(2.75), Inches(3.72), Inches(1.0))
        box(s, x + Inches(0.22), Inches(2.87), Inches(3.3), Inches(0.4), t,
            size=13.5, colour=ACCENT, bold=True)
        box(s, x + Inches(0.22), Inches(3.25), Inches(3.32), Inches(0.5), d,
            size=11, colour=MUTED, space=0)
    box(s, Inches(0.85), Inches(4.0), Inches(11.65), Inches(0.5),
        "The honest part: most municipal CCTV was installed to watch for theft, not to read "
        "plates.", size=14, colour=INK, space=0)
    card(s, Inches(0.85), Inches(4.6), Inches(11.65), Inches(1.05))
    box(s, Inches(1.15), Inches(4.75), Inches(11.1), Inches(0.8),
        "It is mounted high and angled wide. We expect only about **40% of any existing "
        "estate** to be usable for enforcement without repositioning - and we survey it and "
        "say so **before** anyone signs, not after.", size=13.5, colour=INK, space=0)
    note(s, "Where a camera will not do the job, we supply one. That is revenue, not a "
            "problem.", y=5.95, size=14)

    # ------------------------------------------------------------ 11  edge
    s = S(); bg(s)
    header(s, "Why the thinking happens on the pole",
           "STEP 2  -  THE ARITHMETIC THAT DECIDES EVERYTHING")
    table(s, [
        ["Approach", "Network load for 1,000 cameras", "Verdict"],
        ["Send every video stream to a central server",
         "1,000 x 2 Mbps = 2 Gbps sustained, ~21 TB per day", "Impossible"],
        ["Detect at the edge, send EVENTS only",
         "A challan is ~250 KB. ~50 GB/day, bursty - about 5 Mbps average",
         "Runs on what exists"],
    ], widths=[0.31, 0.48, 0.21], top=2.1, rh=0.66, size=11.5)
    box(s, Inches(0.85), Inches(4.2), Inches(11.65), Inches(0.6),
        "A 400x reduction in bandwidth.", size=30, colour=ACCENT, bold=True)
    box(s, Inches(0.85), Inches(4.9), Inches(11.65), Inches(0.5),
        "That is the difference between deployable on Sri Lankan municipal networks, "
        "and not deployable at all.", size=15, colour=INK, space=0)
    for i, (t, d) in enumerate([
            ("Runs on a small computer",
             "CPU is enough today. No GPU cluster, no data centre."),
            ("Survives the link dropping",
             "Events queue locally on a UPS and sync when it returns."),
            ("We keep the frame, not the stream",
             "90-day retention, then automatic deletion.")]):
        x = Inches(0.85 + i * 3.92)
        card(s, x, Inches(5.5), Inches(3.72), Inches(1.15))
        box(s, x + Inches(0.2), Inches(5.63), Inches(3.4), Inches(0.45), t,
            size=12.5, colour=ACCENT, bold=True)
        box(s, x + Inches(0.2), Inches(6.05), Inches(3.4), Inches(0.6), d,
            size=11, colour=MUTED, space=0)

    # ------------------------------------------------------------ 12  the AI
    s = S(); bg(s)
    header(s, "What the AI does with one frame", "STEPS 3 TO 5  -  THE DETECTION FLOW")
    rows(s, [
        ("Detect", "One pass of the model finds vehicles, people, helmets, phones and "
                   "number plates."),
        ("Track", "The same motorcycle keeps ID #47 across frames, so a rule can watch it "
                  "over time instead of judging a snapshot."),
        ("Associate", "Is this person actually ON that motorcycle? Overlap, centre and foot "
                      "position are all checked - so a pedestrian is never fined."),
        ("Persist", "Three to four frames must agree before any accusation. One bad frame "
                    "is never enough."),
        ("Measure", "Is it genuinely moving? We use net displacement with camera shake "
                    "removed, so a parked bike can never be fined."),
        ("Refuse", "If the camera is not calibrated, the rules that need geometry do not "
                   "fire at all. We deleted our own guess rather than ship it."),
    ], top=2.05, gap=0.79, hgt=0.68, label_w=2.0, size=15)
    note(s, "Most of our engineering went into the checks that stop it firing - not the "
            "ones that make it fire. In an evidence product, that is the whole job.",
         y=6.9, size=14)

    # ------------------------------------------------------------ 13  data kinds
    s = S(); bg(s)
    header(s, "Where the zone data comes from",
           "STEP 4  -  THE QUESTION YOU PRESSED HARDEST ON")
    card(s, Inches(0.85), Inches(2.05), Inches(11.65), Inches(1.0))
    box(s, Inches(1.15), Inches(2.2), Inches(11.1), Inches(0.8),
        "Today a human opens each clip, clicks four corners of the road and types the "
        "real distance. That is fine for five cameras. It is impossible for five thousand.",
        size=14.5, colour=INK, space=0)
    box(s, Inches(0.85), Inches(3.35), Inches(11.65), Inches(0.5),
        "Our demo hides something: there are TWO kinds of data, and they scale completely "
        "differently.", size=16, colour=ACCENT, bold=True)
    table(s, [
        ["Kind of data", "Example", "Per camera?", "Cost"],
        ["Camera geometry - where the ground plane sits in THIS image",
         "4-point homography, metres per pixel", "Yes. Unavoidable.", "Minutes, once"],
        ["World knowledge - what the law says about this piece of road",
         "Speed limit, no-parking, one-way, stop line", "No. It belongs to the map.",
         "Once, nationally"],
    ], widths=[0.36, 0.28, 0.22, 0.14], top=4.1, rh=0.78, size=11)
    note(s, "We were treating both as per-camera work. That made the job look a hundred "
            "times bigger than it is.", y=6.1, size=14)

    # ------------------------------------------------------------ 14  geo-register
    s = S(); bg(s)
    header(s, "Geo-register each camera once", "FOUR MINUTES, THEN NEVER AGAIN")
    box(s, Inches(0.85), Inches(2.0), Inches(11.65), Inches(0.6),
        "The operator drags four points on the camera image and the matching four on a "
        "satellite map. That one act then delivers, automatically:", size=15,
        colour=INK, space=0)
    for i, (what, why) in enumerate([
            ("Metres per pixel", "speed measurement"),
            ("Stop-line position", "red-light rule"),
            ("Legal direction of travel", "wrong-way rule"),
            ("No-parking polygons", "illegal parking"),
            ("School zones, bus halts", "context rules"),
            ("Crossings, junction shape", "future rules, free")]):
        x = Inches(0.85 + (i % 3) * 3.92)
        y = Inches(2.75 + (i // 3) * 1.05)
        card(s, x, y, Inches(3.72), Inches(0.9))
        box(s, x + Inches(0.2), y + Inches(0.1), Inches(3.4), Inches(0.4), what,
            size=13, colour=ACCENT, bold=True)
        box(s, x + Inches(0.2), y + Inches(0.46), Inches(3.4), Inches(0.4), "-> " + why,
            size=11, colour=MUTED, space=0)
    for i, (v, l) in enumerate([("4 min", "per camera, once"),
                                ("67 hrs", "for 1,000 cameras"),
                                ("2 people", "2 weeks, a province"),
                                ("Auto", "re-flags if moved")]):
        stat(s, Inches(0.85 + i * 2.95), Inches(5.05), v, l, colour=ACCENT,
             w=Inches(2.75))
    note(s, "If a camera is knocked or moved, the background changes - we detect it, stop "
            "enforcing the geometry rules, and ask to be recalibrated.", y=6.6, size=13.5)

    # ------------------------------------------------------------ 15  import world
    s = S(); bg(s)
    header(s, "Import the world once, nationally",
           "WE DO NOT CREATE THIS DATA - IT ALREADY EXISTS")
    table(s, [
        ["Source", "What it gives us", "Who holds it"],
        ["RDA road register", "Road class and centreline geometry", "Government"],
        ["Survey Dept / provincial GIS",
         "Built-up-area boundaries = the legal 50 km/h limit", "Government"],
        ["Gazetted traffic schemes", "No-parking, one-way, bus halts", "Municipal councils"],
        ["OpenStreetMap", "maxspeed, oneway, crossings, schools - urban coverage today",
         "Free"],
        ["Motor Traffic Act defaults",
         "Where nothing is known: the statutory limit for that road class", "The law itself"],
        ["Learned proposals",
         "'95% of vehicles here do 38-46 km/h - is the limit right?'", "Our own data"],
    ], widths=[0.26, 0.56, 0.18], top=2.05, rh=0.56, size=11)
    note(s, "The councils' schemes are paper and PDF today. Digitising Colombo's is a few "
            "weeks for one person - once, and then never again.", y=6.15, size=13.5)
    note(s, "Learned rules are PROPOSED for a human to approve. They are never auto-enforced.",
         y=6.7, size=13.5, colour=AMBER)

    # ------------------------------------------------------------ 16  day one
    s = S(); bg(s)
    header(s, "Six of ten rules need no map at all",
           "WHY THE DATA PROBLEM DOES NOT BLOCK DAY ONE")
    table(s, [
        ["Rule", "What it needs from the world", "Live on day one?"],
        ["No Helmet", "nothing", "YES"],
        ["Triple Riding", "nothing", "YES"],
        ["Wheelie / stunt riding", "nothing", "YES"],
        ["Mobile Phone Use", "nothing", "YES"],
        ["No Rest Break", "nothing - time only", "YES"],
        ["No Seatbelt", "nothing (needs a working detector)", "model pending"],
        ["Wrong Way", "legal direction of travel", "after geo-registration"],
        ["Over Speeding", "ground plane + posted limit", "after geo-registration"],
        ["Illegal Parking", "no-parking polygon", "after zone import"],
        ["Red Light Jump", "stop line + signal state", "after geo-registration"],
    ], widths=[0.28, 0.46, 0.26], top=1.95, rh=0.40, size=11)
    note(s, "We deploy immediately and enrich over weeks. The customer sees value in week "
            "one, not month six.", y=6.5, size=14)

    # ############################################################ PART 3
    divider(S(), 3, "Money",
            "What it costs, what people pay, and why those numbers are right for Sri Lanka")

    # ------------------------------------------------------------ 18  price list
    s = S(); bg(s)
    header(s, "How you subscribe or licence it", "THE PRICE LIST")
    table(s, [
        ["Tier", "Who it is for", "Price"],
        ["Audit", "First contact - a council or division. 30-day report, no fines issued",
         "Free"],
        ["Enforce", "Per camera. All rules, dashboard, officer review queue",
         "LKR 3,500 / camera / month"],
        ["Enforce + ANPR", "Adds plate recognition, owner lookup, PDF challan",
         "LKR 5,500 / camera / month"],
        ["Fleet", "Bus, logistics, plantation, courier. Driver scorecards",
         "LKR 25,000 / month (25 vehicles)"],
        ["Provincial Licence", "Ministry / Police HQ. Unlimited cameras, ON-PREMISE, SLA",
         "LKR 12,000,000 / province / year"],
        ["Insight", "Insurers, RDA, researchers. Anonymised blackspot analytics",
         "LKR 2,000,000 / year"],
        ["Site Activation", "One-time: survey, geo-registration, zone import, training",
         "LKR 45,000 / camera"],
        ["Edge appliance", "One-time: hardware at cost + 15%, 4 cameras, 3-yr warranty",
         "LKR 138,000 / box"],
    ], widths=[0.19, 0.55, 0.26], top=2.0, rh=0.47, size=11)
    note(s, "A police force will not put criminal evidence in someone else's cloud - and "
            "should not. The on-premise licence is our highest-margin product, not a "
            "compromise.", y=6.45, size=13.5)

    # ------------------------------------------------------------ 19  how we priced it
    s = S(); bg(s)
    header(s, "Five rules we used to set the price", "HOW THE PRICING MODEL WORKS")
    rows(s, [
        ("Anchor on the alternative",
         "Not on our cost. Three constables cover one point 24/7 for ~LKR 300,000 a month. "
         "We charge LKR 5,750."),
        ("Operating cost, not capital",
         "A council can approve a monthly maintenance line. A purchase needs a tender. So we "
         "offer LKR 5,000/camera/month with NO upfront fee, setup rolled in over 3 years."),
        ("Start at one camera",
         "No minimum order. A council can try four cameras without convening a committee."),
        ("Publish volume tiers now",
         "So we are never negotiated down from zero: 100+ cameras LKR 3,000; 500+ LKR 2,600; "
         "2,000+ LKR 2,200."),
        ("Charge fleets per vehicle",
         "A fleet owner thinks in vehicles, not cameras. LKR 1,000 per vehicle per month is a "
         "number they can check against their insurance bill."),
    ], top=2.05, gap=0.93, hgt=0.82, label_w=3.3, size=14)
    note(s, "Free at the top of the ladder, cheap to start, cheaper at scale, and never a "
            "capital purchase. That is the whole model.", y=6.85, size=14)

    # ------------------------------------------------------------ 20  right for SL?
    s = S(); bg(s)
    header(s, "Is this the right price for Sri Lanka?", "THE AFFORDABILITY TEST")
    table(s, [
        ["Buyer", "What they pay us", "Compared with", "Verdict"],
        ["Council, 40 cameras", "LKR 1.68M / year",
         "Less than 2 constables cost for a year", "Fits an existing budget line"],
        ["Police HQ, one province", "LKR 12M / year",
         "A rounding error against the Police budget", "Approvable without a new vote"],
        ["Fleet, 25 vehicles", "LKR 1,000 / vehicle / month",
         "About 10% of a commercial motor premium", "Pays for itself on a 10% discount"],
        ["Versus imported software", "~USD 140 / camera / year",
         "Imported ANPR licences run USD 200-1,000", "We are 3-7x cheaper"],
    ], widths=[0.21, 0.22, 0.32, 0.25], top=2.05, rh=0.62, size=10.5)
    card(s, Inches(0.85), Inches(4.75), Inches(11.65), Inches(1.5))
    box(s, Inches(1.15), Inches(4.92), Inches(11.1), Inches(0.45),
        "The real pricing risk is not the price. It is getting paid.", size=17,
        colour=AMBER, bold=True)
    box(s, Inches(1.15), Inches(5.42), Inches(11.1), Inches(0.8),
        "Government in Sri Lanka pays in 90 to 180 days, and that is what kills small "
        "vendors. So: 25% of activation invoiced up front, the provincial licence billed "
        "quarterly in advance, and private-sector revenue held as working capital.",
        size=13, colour=INK, space=0)
    note(s, "We would rather price it so a council can say yes this quarter than price it "
            "so it looks impressive on a slide.", y=6.45, size=14)

    # ------------------------------------------------------------ 21  no fine share
    s = S(); bg(s)
    header(s, "What we will NOT do: take a share of the fines",
           "THE OBVIOUS ANSWER, AND WHY IT IS WRONG")
    for i, (t, d) in enumerate([
        ("It is probably not legal",
         "Fines under the Motor Traffic Act are public revenue. Assigning a private slice "
         "of them is a constitutional and audit problem."),
        ("It creates exactly the wrong incentive",
         "A company paid per fine wants more fines. Our goal is FEWER violations - under "
         "revenue share, success would bankrupt us."),
        ("It is politically fatal",
         "'Startup profits from fining Sri Lankan motorists' is a headline that ends the "
         "company in a week.")]):
        y = Inches(2.15 + i * 1.35)
        card(s, Inches(0.85), y, Inches(11.65), Inches(1.2))
        box(s, Inches(1.15), y + Inches(0.14), Inches(11.1), Inches(0.45), t,
            size=17, colour=AMBER, bold=True)
        box(s, Inches(1.15), y + Inches(0.6), Inches(11.0), Inches(0.6), d,
            size=13, colour=INK, space=0)
    note(s, "Subscription pays us for coverage and deterrence - which is the outcome the "
            "customer actually wants.", y=6.4)

    # ------------------------------------------------------------ 22  constable
    s = S(); bg(s)
    header(s, "Where LKR 3,500 comes from", "PRICED AGAINST A CONSTABLE")
    table(s, [
        ["", "One traffic constable", "One camera on our platform"],
        ["Cost per month, all-in", "~LKR 100,000",
         "LKR 3,500 + LKR 1,250 amortised setup"],
        ["Hours covered per day", "8 - one shift", "24"],
        ["To cover one point 24/7", "3 officers = ~LKR 300,000 / month",
         "LKR 5,750 / month"],
        ["Rules watched at once", "2-3 realistically", "10"],
        ["Photographic evidence", "No", "Every time"],
    ], widths=[0.28, 0.34, 0.38], top=2.1, rh=0.50, size=11.5)
    box(s, Inches(0.85), Inches(5.25), Inches(11.65), Inches(0.7),
        "Roughly 50x cheaper per point-hour of coverage.", size=30, colour=ACCENT,
        bold=True)
    note(s, "And we are not replacing officers. We are telling them where to stand. That "
            "framing matters when the buyer has a union and a manpower shortage.",
         y=6.1, colour=INK, size=15, bold=False)
    note(s, "Gross margin at LKR 3,500: about 74%.", y=6.75, size=13, colour=MUTED,
         bold=False)

    # ------------------------------------------------------------ 23  province P&L
    s = S(); bg(s)
    header(s, "One province, three years", "WESTERN PROVINCE  -  MODELLED")
    table(s, [
        ["", "Year 1  -  prove it", "Year 2  -  own it", "Year 3  -  expand"],
        ["Cameras live", "80", "600", "2,200 (4 provinces)"],
        ["Subscriptions", "2.9M", "32.4M", "118.8M"],
        ["Site activation + hardware", "6.4M", "23.4M", "72.0M"],
        ["Provincial licences", "-", "12.0M", "48.0M"],
        ["Fleet + insight", "1.6M", "14.0M", "63.0M"],
        ["TOTAL REVENUE", "10.8M  (~$36k)", "81.8M  (~$273k)", "301.8M  (~$1.0M)"],
        ["Total cost", "21.6M", "64.0M", "-"],
        ["NET", "-10.8M", "+17.8M  BREAKEVEN", "-"],
    ], widths=[0.30, 0.235, 0.235, 0.23], top=2.0, rh=0.46, size=11.5)
    note(s, "All figures LKR. Assumed LKR 300 = USD 1. These are our own modelled "
            "assumptions with the working shown - not published facts.", y=6.25, size=13,
         colour=MUTED, bold=False)
    note(s, "Year 1's deficit is the ask: LKR 15M (~USD 50,000) for 18 months of runway.",
         y=6.72, size=15)

    # ------------------------------------------------------------ 24  ceiling
    s = S(); bg(s)
    header(s, "The honest ceiling - and why that is fine",
           "WE WILL SAY IT BEFORE YOU DO")
    card(s, Inches(0.85), Inches(2.05), Inches(11.65), Inches(1.35))
    box(s, Inches(1.15), Inches(2.25), Inches(11.1), Inches(1.0),
        "Fully saturated, Sri Lanka is about 5,000 enforcement cameras, 9 provincial "
        "licences and a few hundred fleets. That is roughly **LKR 480M - about USD 1.6M "
        "of annual recurring revenue.**", size=15.5, colour=INK, space=0)
    box(s, Inches(0.85), Inches(3.65), Inches(11.65), Inches(0.6),
        "Sri Lanka alone is not a big enough market to build a big company in.",
        size=24, colour=AMBER, bold=True)
    box(s, Inches(0.85), Inches(4.45), Inches(11.65), Inches(0.6),
        "It is exactly the right market to prove one in.", size=24, colour=ACCENT, bold=True)
    box(s, Inches(0.85), Inches(5.3), Inches(11.65), Inches(1.5), [
        "The product transfers unchanged to Bangladesh, Nepal, Pakistan, Vietnam, Kenya, "
        "Nigeria, Tanzania.",
        "Same traffic mix - motorcycles, three-wheelers, mixed lanes - which is precisely "
        "what Western products are NOT built for.",
        "Worse death rates, and donor money already attached to fixing them.",
        "**Sri Lanka is the reference customer, not the whole business.**"],
        size=14, colour=INK, space=9)

    # ############################################################ PART 4
    divider(S(), 4, "Implementation",
            "How we launch it, what will go wrong, and what we need from you")

    # ------------------------------------------------------------ 26  roadmap
    s = S(); bg(s)
    header(s, "How it actually reaches a road", "IMPLEMENTATION ROADMAP")
    table(s, [
        ["Phase", "Months", "What we do", "Done when"],
        ["0  Harden", "0-3", "Incorporate. Build a labelled Sri Lankan ground-truth set "
                             "and publish real accuracy. Officer review queue. PDPA assessment.",
         "We can state an accuracy number and defend it"],
        ["1  Audit", "2-5", "Free 30-day audit, 2 junctions, 1 council. No fines.",
         "A signed report with a council logo on it"],
        ["2  Champion", "4-7", "Take the report to NCRS and Police Traffic HQ.",
         "A letter of support in hand"],
        ["3  Pilot", "6-12", "20-40 cameras, one division, every challan human-approved.",
         "Disputes upheld near zero"],
        ["4  Commercial", "3-12", "In parallel: 8 fleets, 2 campuses or ports.",
         "Revenue covers half the burn"],
        ["5  Province", "12-24", "Western Province at 600 cameras. Integrate DMT and e-payment.",
         "Breakeven"],
        ["6  Scale", "24-36", "3 more provinces. First regional pilot.", "USD 1M run rate"],
    ], widths=[0.14, 0.09, 0.46, 0.31], top=1.92, rh=0.57, size=10.5)
    note(s, "The longest-lead item is the Department of Motor Traffic data-sharing agreement "
            "- plate to registered owner. We start it in phase 2, not phase 5.", y=6.7,
         size=13)

    # ------------------------------------------------------------ 27  challenges
    s = S(); bg(s)
    header(s, "What will go wrong, and what we do about it", "CHALLENGES")
    table(s, [
        ["Challenge", "What we do about it"],
        ["Procurement takes two years",
         "Private revenue funds us; sub-contract under a registered integrator"],
        ["Is AI evidence admissible?",
         "Officer approves every fine; audit trail; calibration certificate per camera"],
        ["Privacy backlash / PDPA 2022",
         "No face recognition, ever. Plates only. 90-day retention, then auto-purge"],
        ["'Big Brother' media story",
         "Lead with lives, not fines. Publish violation-rate falls, not revenue collected"],
        ["Existing cameras are the wrong quality",
         "Survey first, state the usable percentage up front, sell our own where needed"],
        ["Government pays 90-180 days late",
         "Bill activation up front, licences quarterly in advance, private revenue as buffer"],
        ["A false positive fines an innocent person",
         "Human approval gate, precision over recall, one-click dispute, publish our error rate"],
        ["Power cuts and network drops",
         "Edge box queues locally on a UPS and syncs when the link returns"],
        ["'Make this challan disappear'",
         "Append-only audit log; every view and cancellation attributed to a named officer"],
        ["Officers resist - it threatens discretion",
         "Position it as deployment intelligence: it tells them where to stand"],
    ], widths=[0.36, 0.64], top=1.92, rh=0.45, size=11)
    note(s, "None of these are new. Every country that automated enforcement hit all of them.",
         y=6.75, size=13.5, colour=MUTED, bold=False)

    # ------------------------------------------------------------ 28  SWOT
    s = S(); bg(s)
    header(s, "SWOT", "AN HONEST ONE")
    for i, (t, colour, items) in enumerate([
        ("STRENGTHS", GREEN, [
            "It exists - 10 rules on live video, 85 tests, our own evidence",
            "Built for Sri Lankan traffic: bikes, three-wheelers, mixed lanes",
            "It refuses to guess - no calibration, no speed reading",
            "CPU-capable and edge-first: cheap to run",
            "Local team, local cost base, same-afternoon support"]),
        ("WEAKNESSES", AMBER, [
            "No revenue, no reference customer, no registered company",
            "No measured accuracy figure yet - a court will ask",
            "Seatbelt rule disabled: the model we got is unusable",
            "Three part-time students; no procurement or legal experience",
            "We depend on camera quality we do not control"]),
        ("OPPORTUNITIES", ACCENT, [
            "~3,000 deaths a year; enforcement is still manual and spot-based",
            "Existing municipal CCTV is a sunk asset we can activate",
            "PDPA 2022 compliance is a moat foreign vendors handle badly",
            "Donor funding needs no police budget line",
            "Fleet and insurance markets need no government sale at all"]),
        ("THREATS", INK, [
            "Hikvision / Dahua / Huawei bundling ANPR free with hardware",
            "Tenders written around an incumbent, or never issued",
            "Political change resetting our sponsor",
            "One high-profile false positive destroying trust",
            "We are students - graduation is a continuity risk"])]):
        x = Inches(0.85 + (i % 2) * 6.0)
        y = Inches(1.95 + (i // 2) * 2.6)
        card(s, x, y, Inches(5.65), Inches(2.45))
        box(s, x + Inches(0.25), y + Inches(0.14), Inches(5.1), Inches(0.4), t,
            size=13.5, colour=colour, bold=True)
        box(s, x + Inches(0.25), y + Inches(0.58), Inches(5.15), Inches(1.8), items,
            size=10.5, colour=INK, space=6)

    # ------------------------------------------------------------ 29  moat
    s = S(); bg(s)
    header(s, "The model is not the moat", "COMPETITION AND DEFENSIBILITY")
    box(s, Inches(0.85), Inches(1.95), Inches(11.65), Inches(0.7),
        "Anyone can download YOLO. Our biggest threat is a global vendor giving the software "
        "away to win the camera contract - so we do not compete on price. We run on top of "
        "their hardware and sell the layer they cannot.", size=14.5, colour=INK, space=0)
    rows(s, [
        ("The calibrated camera estate",
         "1,000 geo-registered cameras is months of field work a competitor must repeat"),
        ("Workflow lock-in",
         "Once challans flow through our review queue, switching costs are retraining"),
        ("A labelled Sri Lankan dataset",
         "Plates, helmets, three-wheelers, night, monsoon - no foreign vendor has this"),
        ("Legal precedent",
         "The first vendor whose challan survives a court challenge owns the category"),
        ("Cost structure",
         "A local team on an LKR cost base undercuts anyone importing engineers"),
    ], top=2.75, gap=0.85, hgt=0.72, label_w=3.9, size=13.5)

    # ------------------------------------------------------------ 30  the ask
    s = S(); bg(s)
    header(s, "What we are asking for", "THE ASK")
    for i, (v, l) in enumerate([("LKR 15M", "~USD 50,000 seed"),
                                ("18 months", "of runway"),
                                ("80 cameras", "live by month 12"),
                                ("10", "paying customers")]):
        stat(s, Inches(0.85 + i * 2.95), Inches(2.05), v, l, colour=ACCENT, w=Inches(2.75))
    box(s, Inches(0.85), Inches(3.65), Inches(11.65), Inches(0.5),
        "60% salaries  |  15% pilot hardware  |  10% legal and PDPA compliance  |  "
        "10% field operations  |  5% buffer", size=13, colour=MUTED, space=0)
    box(s, Inches(0.85), Inches(4.3), Inches(11.65), Inches(0.5),
        "And four things that are not money:", size=17, colour=INK, bold=True)
    box(s, Inches(1.1), Inches(4.9), Inches(11.4), Inches(2.0), [
        "**1.**  One letter of support from the National Council for Road Safety or Police Traffic HQ.",
        "**2.**  One municipal council willing to host a free 30-day audit.",
        "**3.**  An introduction to a government IT integrator with existing vendor registration.",
        "**4.**  Access to a labelled Sri Lankan traffic dataset - or permission to build one."],
        size=14.5, colour=INK, space=11)

    # ------------------------------------------------------------ 31  close
    s = S(); bg(s)
    faded_bg(s, "06_junction_tracking_0.jpg", opacity=0.13)
    accent_bar(s)
    box(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.2),
        "We already built the hard part.", size=42, colour=INK, bold=True)
    box(s, Inches(0.9), Inches(3.6), Inches(11.5), Inches(0.7),
        "Now we are asking for the boring part: one council, one letter, one pilot.",
        size=20, colour=ACCENT)
    box(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.6),
        "Thank you. We are happy to take your questions.", size=17, colour=INK)
    box(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.6),
        "Team Three Hacks  |  github.com/Jeba-Jebarsan/AI-SMART-TRAFFIC-INTELLIGENCE",
        size=13, colour=MUTED)

    names = (["PRESENTATION_ROUND2.pptx"]
             + [f"PRESENTATION_ROUND2_v{i}.pptx" for i in range(2, 9)])
    for i, name in enumerate(names):
        out = os.path.join(ROOT, "docs", name)
        try:
            prs.save(out)
        except PermissionError:
            continue
        if i:
            print(f"NOTE: {names[0]} is open in PowerPoint - wrote {name} instead. "
                  f"Close the deck and re-run to write the real file.")
        print("wrote", out, "|", len(prs.slides._sldIdLst), "slides")
        return
    raise SystemExit("every PRESENTATION_ROUND2*.pptx is locked - close PowerPoint")


if __name__ == "__main__":
    main()
