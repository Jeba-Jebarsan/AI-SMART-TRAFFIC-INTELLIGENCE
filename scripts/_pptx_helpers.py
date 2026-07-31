"""Shared drawing helpers for the PowerPoint deck.

Light professional theme: white page, near-black text, one teal accent. A dark
deck reads well on a laptop but loses contrast when a judge is watching a
compressed Google Meet stream on their own screen, and it prints badly if the
panel wants a handout. White is the safer room-independent choice.
"""
import os

from PIL import Image as PILImage
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
EV = os.path.join(ROOT, "docs", "DEMO_EVIDENCE")

W, H = Inches(13.333), Inches(7.5)          # 16:9

BG = RGBColor(0xFF, 0xFF, 0xFF)             # page
INK = RGBColor(0x0F, 0x17, 0x2A)            # primary text
MUTED = RGBColor(0x5B, 0x6B, 0x84)          # secondary text
ACCENT = RGBColor(0x0E, 0x74, 0x90)         # teal - headings, rules
PANEL = RGBColor(0xF1, 0xF5, 0xF9)          # card fill
PANEL_LINE = RGBColor(0xCB, 0xD5, 0xE1)     # card border
AMBER = RGBColor(0xB4, 0x53, 0x09)          # emphasis / key takeaway
GREEN = RGBColor(0x15, 0x80, 0x3D)          # verified / positive


def bg(slide, colour=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = colour


def box(slide, x, y, w, h, text, size=18, colour=INK, bold=False,
        align=PP_ALIGN.LEFT, font="Segoe UI", space=6):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    lines = text if isinstance(text, list) else [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        # inline **bold** segments
        for j, seg in enumerate(str(line).split("**")):
            if not seg:
                continue
            r = p.add_run()
            r.text = seg
            r.font.size = Pt(size)
            r.font.name = font
            r.font.color.rgb = colour
            r.font.bold = bold or (j % 2 == 1)
    return tb


def accent_bar(slide):
    s = slide.shapes.add_shape(1, 0, 0, W, Inches(0.09))
    s.fill.solid()
    s.fill.fore_color.rgb = ACCENT
    s.line.fill.background()
    s.shadow.inherit = False


def card(slide, x, y, w, h):
    """A filled panel with a hairline border - the deck's only container."""
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = PANEL
    s.line.color.rgb = PANEL_LINE
    s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def picture(slide, name, x, y, max_w, max_h):
    """Place an image scaled to fit the given box, centred horizontally."""
    path = os.path.join(EV, name)
    if not os.path.exists(path):
        return None
    iw, ih = PILImage.open(path).size
    scale = min(max_w / iw, max_h / ih)
    w, h = int(iw * scale), int(ih * scale)
    return slide.shapes.add_picture(path, x + int((max_w - w) / 2), y,
                                    width=w, height=h)


def faded_bg(slide, name, opacity=0.16, cache_dir=None):
    """Full-bleed background from one of OUR OWN evidence frames, faded to near-white.

    Deliberately not a stock photo: the title of a deck about a working system
    should be the system's own output. Faded this far it reads as texture, and
    dark title text stays legible over it.

    The fade is baked into a cached JPEG rather than done with PowerPoint
    transparency, because picture transparency needs raw XML in python-pptx and
    renders inconsistently across PowerPoint versions.
    """
    src = os.path.join(EV, name)
    if not os.path.exists(src):
        return None
    cache_dir = cache_dir or os.path.join(ROOT, "docs", "DEMO_EVIDENCE", "_derived")
    os.makedirs(cache_dir, exist_ok=True)
    out = os.path.join(cache_dir, f"bg_{int(opacity * 100)}_{name}")

    if not os.path.exists(out) or os.path.getmtime(out) < os.path.getmtime(src):
        img = PILImage.open(src).convert("RGB")
        target = 16 / 9
        w, h = img.size
        if w / h > target:                      # too wide -> crop sides
            nw = int(h * target)
            img = img.crop(((w - nw) // 2, 0, (w - nw) // 2 + nw, h))
        else:                                   # too tall -> crop top/bottom
            nh = int(w / target)
            img = img.crop((0, (h - nh) // 2, w, (h - nh) // 2 + nh))
        img = img.resize((1920, 1080), PILImage.LANCZOS)
        white = PILImage.new("RGB", img.size, (255, 255, 255))
        PILImage.blend(white, img, opacity).save(out, quality=88)

    pic = slide.shapes.add_picture(out, 0, 0, width=W, height=H)
    # Send behind everything else added afterwards.
    slide.shapes._spTree.remove(pic._element)
    slide.shapes._spTree.insert(2, pic._element)
    return pic


def header(slide, title, kicker=None):
    accent_bar(slide)
    if kicker:
        box(slide, Inches(0.7), Inches(0.42), Inches(11), Inches(0.4),
            kicker, size=13, colour=ACCENT, bold=True)
    box(slide, Inches(0.7), Inches(0.78), Inches(12), Inches(1.0),
        title, size=34, colour=INK, bold=True)


def stat(slide, x, y, value, label, colour=AMBER, w=Inches(2.9)):
    card(slide, x, y, w, Inches(1.35))
    box(slide, x, y + Inches(0.14), w, Inches(0.6), value, size=30,
        colour=colour, bold=True, align=PP_ALIGN.CENTER)
    box(slide, x, y + Inches(0.78), w, Inches(0.5), label, size=11.5,
        colour=MUTED, align=PP_ALIGN.CENTER)
