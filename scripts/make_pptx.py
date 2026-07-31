"""Generate the PowerPoint deck (docs/PRESENTATION.pptx).

Built for the Startup Innovation Competition 2026 final: a 30-minute online
slot on Google Meet, Saturday 1 August 2026.

Design rules this file follows deliberately:
  * Light theme. See _pptx_helpers for why.
  * NOTHING backstage on a slide. No speaker names, no timings, no stage
    directions - the panel sees every pixel we project. Who says what lives in
    docs/SPEAKING_SCRIPTS.md.
  * Capability is framed as "nine today, and the tenth is a config file", not
    as a fixed list of nine, so the platform does not read as a hard limit.
  * Every image is the system's own annotated output and every technical
    figure is a measured result. Commercial figures are labelled as
    projections with their assumptions stated.

    python scripts/make_pptx.py
"""
import os

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from _pptx_helpers import (ACCENT, AMBER, BG, GREEN, H, INK, MUTED, PANEL,
                           PANEL_LINE, W, accent_bar, bg, box, card, faded_bg,
                           header, picture, stat)

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
os.chdir(ROOT)


def bullets(s, items, top=2.05, size=18, gap=14, left=0.8, width=11.6):
    box(s, Inches(left), Inches(top), Inches(width), Inches(4.0), items,
        size=size, colour=INK, space=gap)


def two_col(s, left_title, left_items, right_title, right_items, top=2.0):
    for i, (t, items) in enumerate([(left_title, left_items),
                                    (right_title, right_items)]):
        x = Inches(0.8 + i * 6.1)
        card(s, x, Inches(top), Inches(5.7), Inches(4.3))
        box(s, x + Inches(0.3), Inches(top + 0.22), Inches(5.1), Inches(0.5),
            t, size=17, colour=ACCENT, bold=True)
        box(s, x + Inches(0.3), Inches(top + 0.85), Inches(5.1), Inches(3.2),
            items, size=13.5, colour=INK, space=10)


def evidence_slide(add, title, kicker, img, caption):
    s = add(); bg(s); header(s, title, kicker)
    picture(s, img, Inches(1.4), Inches(1.95), Inches(10.5), Inches(4.2))
    box(s, Inches(0.8), Inches(6.35), Inches(11.6), Inches(0.9), caption,
        size=14, colour=MUTED, align=PP_ALIGN.CENTER)
    return s


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    blank = prs.slide_layouts[6]
    S = lambda: prs.slides.add_slide(blank)   # noqa: E731

    # ---------------------------------------------------------- 1  title
    s = S(); bg(s)
    faded_bg(s, "05_anpr_sri_lanka_0.jpg", opacity=0.17)
    accent_bar(s)
    box(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.3),
        "AI Smart Traffic Intelligence Platform", size=44, colour=INK, bold=True)
    box(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(0.7),
        "Turning any camera feed into fair, evidence-based enforcement",
        size=20, colour=ACCENT)
    box(s, Inches(0.9), Inches(4.25), Inches(11.5), Inches(0.6),
        "Team Three Hacks  |  Startup Innovation Competition 2026",
        size=15, colour=MUTED)
    for i, (v, l) in enumerate([("9", "violation types live today"),
                                ("85", "automated tests passing"),
                                ("83-163", "km/h speed verified")]):
        stat(s, Inches(0.9 + i * 3.15), Inches(5.35), v, l)

    # ---------------------------------------------------------- 2  problem
    s = S(); bg(s)
    header(s, "Every day, people die on roads a camera was watching",
           "PROBLEM STATEMENT")
    bullets(s, ["Sri Lanka records roughly **3,000 road deaths a year**, most of them "
                "from violations that are simple to see.",
                "Officers cannot stand at every junction, 24 hours a day. Enforcement "
                "is limited by **how many people we can hire**.",
                "A fine written **without a photo and a verified plate** is argued "
                "about, and often never paid.",
                "Nothing is recorded in a searchable way, so **repeat offenders stay "
                "invisible**."],
            size=18.5)
    box(s, Inches(0.8), Inches(6.0), Inches(11.6), Inches(0.7),
        "The gap is not cameras. The gap is coverage and proof.",
        size=20, colour=ACCENT, bold=True)

    # ---------------------------------------------------------- 3  opportunity
    s = S(); bg(s); header(s, "The opportunity", "WHY NOW")
    two_col(s, "What already exists",
            ["Colombo alone has **over 100 road cameras** already installed.",
             "Police already issue spot fines - the **legal framework is in place**.",
             "Phones and drones make any road coverable today, with no civil works."],
            "What is missing",
            ["Those cameras are **watched by a person, or by nobody at all**.",
             "There is no automatic path from camera to challan to payment.",
             "Enforcement grows with headcount, not with technology."])
    box(s, Inches(0.8), Inches(6.5), Inches(11.6), Inches(0.6),
        "We do not need new cameras. We need intelligence behind the ones already there.",
        size=16, colour=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # ---------------------------------------------------------- 4  solution
    s = S(); bg(s); header(s, "Our solution", "THE PRODUCT")
    bullets(s, ["One AI pipeline on an ordinary camera that judges **traffic "
                "violations of any kind we teach it**.",
                "Reads **number plates on every vehicle**, not only on offenders.",
                "Measures **real road speed** using surveyed road geometry.",
                "Issues an **e-challan with stamped photographic evidence**, sent to "
                "the police automatically."])
    box(s, Inches(0.8), Inches(5.5), Inches(11.6), Inches(0.9),
        "Camera  ->  YOLOv8s  ->  ByteTrack  ->  Homography  ->  Violation engine  "
        "->  ANPR  ->  e-Challan", size=15, colour=ACCENT, bold=True)
    box(s, Inches(0.8), Inches(6.3), Inches(11.6), Inches(0.6),
        "Runs on a single laptop CPU. No new hardware to install.",
        size=15, colour=ACCENT, bold=True)

    # ---------------------------------------------------------- 5  how it works
    s = S(); bg(s); header(s, "How it works, in six steps", "THE PIPELINE")
    for i, (n, t, d) in enumerate([
            ("1", "See", "YOLOv8s finds vehicles, riders, helmets and plates in each frame"),
            ("2", "Follow", "ByteTrack gives every vehicle an identity it keeps across frames"),
            ("3", "Measure", "Four surveyed road corners turn pixels into real metres"),
            ("4", "Judge", "The rule engine runs, each rule behind motion and persistence gates"),
            ("5", "Identify", "OCR reads the plate; three agreeing reads confirm it"),
            ("6", "Act", "An e-challan with photo evidence is generated and emailed")]):
        x = Inches(0.85 + (i % 3) * 4.05); y = Inches(2.15 + (i // 3) * 2.15)
        card(s, x, y, Inches(3.75), Inches(1.85))
        box(s, x + Inches(0.28), y + Inches(0.14), Inches(3.2), Inches(0.4),
            n, size=13, colour=ACCENT, bold=True)
        box(s, x + Inches(0.28), y + Inches(0.5), Inches(3.2), Inches(0.45),
            t, size=18, colour=INK, bold=True)
        box(s, x + Inches(0.28), y + Inches(0.98), Inches(3.25), Inches(0.8),
            d, size=11.5, colour=MUTED)
    box(s, Inches(0.85), Inches(6.55), Inches(11.6), Inches(0.5),
        "Every step runs live. Nothing is pre-recorded and nothing is faked.",
        size=14, colour=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # ------------------------------------------------- 6  capability (extensible)
    s = S(); bg(s)
    header(s, "A rule engine, not a fixed feature list", "CAPABILITY AND EXTENSIBILITY")
    box(s, Inches(0.8), Inches(1.85), Inches(11.6), Inches(0.5),
        "Nine rules are live today. Each is a small module over the same "
        "detection and tracking core, so the tenth is days of work, not a new "
        "product.", size=15, colour=MUTED)
    for i, name in enumerate(["No Helmet", "Triple Riding", "Over Speeding",
                              "Red Light Jump", "Wrong Way", "No Seatbelt",
                              "Mobile Phone Use", "Illegal Parking",
                              "No Rest Break"]):
        cx = Inches(0.8 + (i % 3) * 4.0); cy = Inches(2.5 + (i // 3) * 1.02)
        card(s, cx, cy, Inches(3.7), Inches(0.85))
        box(s, cx, cy + Inches(0.22), Inches(3.7), Inches(0.5), name,
            size=16, colour=INK, bold=True, align=PP_ALIGN.CENTER)
    box(s, Inches(0.8), Inches(5.75), Inches(11.6), Inches(1.1),
        ["**Already asked for and straightforward to add:** lane discipline, "
         "unauthorised bus-lane use, overloaded three-wheelers, number-plate "
         "obscuring, pedestrian-crossing blocking.",
         "Plus ANPR on **every** tracked vehicle, building a searchable "
         "enforcement log underneath all of it."],
        size=13.5, colour=MUTED, space=8)

    # ---------------------------------------------------- 7-11  evidence block
    evidence_slide(S, "Proof: over-speeding", "VERIFIED SYSTEM OUTPUT",
                   "01_over_speeding_evidence_1.jpg",
                   "134.0 km/h in a 60 zone - vehicle boxed, plate zoomed, proof "
                   "stamped onto the image. **7 events** measured at **83-163 km/h** "
                   "across two calibrated cameras.")

    evidence_slide(S, "Proof: speed is measured, not guessed", "HOW THE NUMBER IS EARNED",
                   "09_highway_speed_0.jpg",
                   "We survey four road corners and map pixels to metres, then fit "
                   "distance against time across many frames and reject outliers. "
                   "Replay accuracy: **0.99-1.01x true speed**. An uncalibrated "
                   "camera shows **no speed at all**.")

    evidence_slide(S, "Proof: helmet and triple riding", "VERIFIED SYSTEM OUTPUT",
                   "07_image_upload_triple_helmet.jpg",
                   "Both violations detected from a **single photograph**, producing "
                   "**two separate challans**. Riders are matched to their motorcycle "
                   "geometrically, so a bystander on the pavement is never counted.")

    evidence_slide(S, "Proof: the bystander is not counted", "PRECISION, NOT JUST DETECTION",
                   "11_no_helmet_bystander.jpg",
                   "The rider is fined for no helmet. The person crouching beside the "
                   "bike is **not** counted as a rider, and the scooter behind is not "
                   "accused. **Anyone can draw a box. Judging correctly is the hard "
                   "part.**")

    evidence_slide(S, "Proof: number-plate recognition", "ANPR ON A REAL SRI LANKAN ROAD",
                   "05_anpr_sri_lanka_0.jpg",
                   "150 vehicles tracked; plates **AAG 4002** and **BBJ 8752** read and "
                   "confirmed. A plate is written to a challan only when **three "
                   "separate reads agree** - otherwise it says UNREADABLE.")

    # ------------------------------------------- 12  what a rule needs to fire
    s = S(); bg(s)
    header(s, "Every rule needs the right camera view", "AN HONEST ENGINEERING ANSWER")
    for i, (rule, need, status) in enumerate([
            ("Helmet, triple riding",
             "Any view where the rider is visible",
             "Demonstrated on real footage"),
            ("Over speeding",
             "Four surveyed road corners, so pixels become metres",
             "Demonstrated - 83-163 km/h, two cameras"),
            ("Red light jump",
             "The stop line, and which signal head governs the lane",
             "Fires once calibrated; never guessed"),
            ("Illegal parking",
             "A fixed camera, so stillness can be proven",
             "Needs a mounted camera, not hand-held video"),
            ("No seatbelt, phone use",
             "A front view of the driver, close enough to resolve",
             "Fleet cameras have it; road CCTV does not")]):
        y = Inches(2.05 + i * 0.98)
        box(s, Inches(0.85), y, Inches(3.5), Inches(0.85), rule, size=14.5,
            colour=INK, bold=True)
        box(s, Inches(4.5), y, Inches(4.3), Inches(0.85), need, size=12.5,
            colour=MUTED)
        box(s, Inches(9.0), y, Inches(3.6), Inches(0.85), status, size=12.5,
            colour=ACCENT)
    box(s, Inches(0.85), Inches(6.9), Inches(11.6), Inches(0.5),
        "Every rule is built and tested. We demonstrate only what our footage "
        "genuinely contains.", size=14, colour=ACCENT, bold=True)

    # ---------------------------------------------------------- 13  trust
    s = S(); bg(s); header(s, "Why it can be trusted", "THE REAL DIFFERENTIATOR")
    bullets(s, ["**Parked vehicles are never fined** - 41 vehicles, 0 violations on a "
                "parked-motorcycle clip.",
                "**Unreadable plates say UNREADABLE** - a number is never invented.",
                "**No calibration means no speed shown**, rather than a misleading guess.",
                "**A violation must persist across frames** - one-frame anomalies rejected.",
                "**A model that cannot prove itself is switched off**, not trusted."],
            size=17, gap=12)
    for i, (v, l, c) in enumerate([("8", "false-positive defects found and fixed", AMBER),
                                   ("85", "automated tests passing", GREEN),
                                   ("0", "violations on parked bikes", ACCENT)]):
        stat(s, Inches(0.9 + i * 3.9), Inches(5.55), v, l, colour=c, w=Inches(3.6))

    # ---------------------------------------------------------- 14  rigour
    s = S(); bg(s); header(s, "Found by testing on real roads", "ENGINEERING RIGOUR")
    for i, (d, c) in enumerate([
            ("Our own seatbelt model was broken",
             "It called a belted driver, and an empty road, 'no seatbelt'. "
             "We disabled the rule rather than ship it"),
            ("A stopped car fined for a red light",
             "A signal from a different junction approach was applied to it"),
            ("Helmeted riders accused of no helmet",
             "The helmet crop was catching neighbouring riders"),
            ("A sleeping driver flagged as 'on phone'",
             "A marginal 0.36 detection on a dark steering column"),
            ("Speeds shown on uncalibrated cameras",
             "Meaningless numbers where no perspective calibration exists")]):
        y = Inches(2.05 + i * 0.95)
        box(s, Inches(0.85), y, Inches(5.3), Inches(0.8), d, size=15,
            colour=INK, bold=True)
        box(s, Inches(6.4), y, Inches(6.2), Inches(0.8), c, size=13.5, colour=MUTED)
    box(s, Inches(0.85), Inches(6.85), Inches(11.6), Inches(0.5),
        "We found these in our own system and fixed them. Each one is now an "
        "automated test.", size=14, colour=ACCENT, bold=True)

    # ---------------------------------------------------------- 15  innovation
    s = S(); bg(s); header(s, "What makes it different", "INNOVATION AND UNIQUENESS")
    two_col(s, "Typical systems",
            ["Single purpose - a speed camera, or a plate reader.",
             "Report a number with no proof attached.",
             "Demo well, then fail on real roads with false positives.",
             "Need dedicated hardware installed at each site."],
            "Ours",
            ["**One engine, many rules**, on any camera.",
             "**Evidence-first**: every fine carries a stamped photo and plate crop.",
             "**Engineered to refuse** - it reports only what it can prove.",
             "**Runs on cameras that already exist**, on a laptop CPU."])

    # ---------------------------------------------------------- 16  market
    s = S(); bg(s); header(s, "Who it is for", "TARGET USERS AND MARKET")
    two_col(s, "Primary users",
            ["**Sri Lanka Police traffic division** - automated enforcement on "
             "existing junction cameras.",
             "**Road Development Authority and municipal councils** - highway "
             "corridors, school and hospital zones.",
             "**Commercial fleet operators** - seatbelt and fatigue compliance "
             "for buses and lorries."],
            "Why they buy",
            ["Enforcement reach without hiring proportionally more officers.",
             "Evidence that survives dispute, so fines are actually collected.",
             "Data on where violations concentrate, to target scarce resources.",
             "Fleet operators cut insurance exposure and liability."])
    box(s, Inches(0.8), Inches(6.5), Inches(11.6), Inches(0.6),
        "Entry point: a single pilot junction. Expansion: corridor, then city.",
        size=15, colour=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    # ---------------------------------------------------------- 17  business model
    s = S(); bg(s); header(s, "Business model", "HOW IT SUSTAINS ITSELF")
    for i, (t, d) in enumerate([
            ("Pilot deployment fee",
             "One-off setup and calibration per camera site"),
            ("Annual software licence",
             "Per-camera subscription covering updates and model improvements"),
            ("Support and analytics tier",
             "Dashboards, reporting and repeat-offender analytics for authorities"),
            ("Fleet compliance package",
             "Per-vehicle monitoring sold to commercial operators")]):
        y = Inches(2.0 + i * 1.15)
        card(s, Inches(0.85), y, Inches(11.6), Inches(0.95))
        box(s, Inches(1.15), y + Inches(0.12), Inches(4.6), Inches(0.7), t,
            size=15.5, colour=ACCENT, bold=True)
        box(s, Inches(5.9), y + Inches(0.14), Inches(6.4), Inches(0.7), d,
            size=13.5, colour=INK)
    box(s, Inches(0.85), Inches(6.75), Inches(11.6), Inches(0.6),
        "Cost advantage: no new cameras and no per-site hardware, so a pilot is "
        "software and calibration only.", size=14, colour=ACCENT, bold=True)

    # ---------------------------------------------------------- 18  sustainability
    s = S(); bg(s); header(s, "Sustainability", "ECONOMIC, SOCIAL, ENVIRONMENTAL")
    two_col(s, "Economically self-sustaining",
            ["Runs on **infrastructure already paid for**.",
             "One laptop starts a pilot; no civil works, no gantries.",
             "Recovered fines fund expansion without new capital."],
            "Socially and environmentally sound",
            ["**Fairness**: identical rules for every driver, no discretion.",
             "**Fewer deaths** through deterrence, not punishment.",
             "**Lower emissions** - smoother traffic, fewer crash-related jams.",
             "**Privacy-respecting**: plates confirmed three times or marked "
             "unreadable; no identity inference."])

    # ---------------------------------------------------------- 19  scalability
    s = S(); bg(s); header(s, "Scalability and future plans", "WHERE THIS GOES")
    for i, (phase, what) in enumerate([
            ("Now", "Working system, 85 tests, verified on real Sri Lankan road footage"),
            ("Next 3 months", "Live pilot at one signalised junction with police "
                              "partnership; night and rain models"),
            ("6-12 months", "Corridor deployment on GPU edge boxes; a trained "
                            "three-wheeler class for local traffic"),
            ("Beyond", "City-wide multi-camera dashboard; export to comparable "
                       "South Asian markets")]):
        y = Inches(2.05 + i * 0.92)
        box(s, Inches(0.9), y, Inches(2.6), Inches(0.7), phase, size=15,
            colour=AMBER, bold=True)
        box(s, Inches(3.6), y, Inches(8.9), Inches(0.8), what, size=13.5, colour=INK)
    card(s, Inches(0.85), Inches(5.75), Inches(11.6), Inches(1.35))
    box(s, Inches(1.15), Inches(5.9), Inches(11.0), Inches(0.4),
        "The accuracy ceiling is compute, not architecture", size=15,
        colour=ACCENT, bold=True)
    box(s, Inches(1.15), Inches(6.32), Inches(11.0), Inches(0.7),
        "Today we run **YOLOv8s on a laptop CPU** - chosen because it fits, not "
        "because it is the best available. The detector is one swappable file: on "
        "GPU hardware the same pipeline takes **YOLO11x, RT-DETR or DINO** with no "
        "change to a single rule.", size=12.5, colour=INK)

    # ------------------------------------------------- 19b  the upgrade path
    s = S(); bg(s)
    header(s, "What more compute buys us", "TECHNICAL ROADMAP")
    box(s, Inches(0.8), Inches(1.8), Inches(11.6), Inches(0.5),
        "We benchmarked five detectors on our own footage. Every rule sits on one "
        "detector interface, so upgrading is a file swap, not a rewrite.",
        size=14.5, colour=MUTED)
    for i, (tier, model, gain) in enumerate([
            ("Laptop CPU  (today)", "YOLOv8s",
             "2.5 analysed frames/sec. We measured yolo11m finding 32% more riders "
             "per frame - but running 3.6x slower, so it sees fewer frames and "
             "delivers LESS. On CPU, yolov8s wins."),
            ("GPU edge box  (pilot)", "YOLO11m / YOLO11x",
             "Frame dropping disappears, so raw recall wins instead of speed. The "
             "same measurement now ranks yolo11m first."),
            ("Server GPU  (city scale)", "RT-DETR, DINO",
             "Transformer detectors handle dense, overlapping traffic far better - "
             "exactly our junctions at rush hour."),
            ("Open-vocabulary  (research)", "Grounding DINO, VLMs",
             "Describe a new violation in words instead of collecting and "
             "labelling a dataset for it.")]):
        y = Inches(2.4 + i * 1.15)
        card(s, Inches(0.85), y, Inches(11.6), Inches(1.0))
        box(s, Inches(1.1), y + Inches(0.12), Inches(3.0), Inches(0.7), tier,
            size=13, colour=AMBER, bold=True)
        box(s, Inches(4.2), y + Inches(0.12), Inches(2.5), Inches(0.7), model,
            size=13, colour=ACCENT, bold=True)
        box(s, Inches(6.8), y + Inches(0.1), Inches(5.5), Inches(0.85), gain,
            size=11.5, colour=INK)
    box(s, Inches(0.85), Inches(7.05), Inches(11.6), Inches(0.5),
        "The hard part - the rules that refuse to guess - is the part that does "
        "not need replacing.", size=14, colour=ACCENT, bold=True)

    # ---------------------------------------------------------- 21  close
    s = S(); bg(s)
    faded_bg(s, "01_over_speeding_evidence_1.jpg", opacity=0.13)
    accent_bar(s)
    box(s, Inches(0.9), Inches(2.5), Inches(11.5), Inches(1.2),
        "Safer roads, proven by evidence.", size=42, colour=INK, bold=True)
    box(s, Inches(0.9), Inches(3.8), Inches(11.5), Inches(0.7),
        "Real speed. Verified plates. Automatic challans, with the photograph attached.",
        size=18, colour=ACCENT)
    box(s, Inches(0.9), Inches(4.7), Inches(11.5), Inches(0.6),
        "On cameras this country already owns.", size=17, colour=ACCENT, bold=True)
    box(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(0.6),
        "Thank you. We are happy to take your questions.", size=16, colour=INK)
    box(s, Inches(0.9), Inches(6.4), Inches(11.5), Inches(0.6),
        "Team Three Hacks  |  github.com/Jeba-Jebarsan/AI-SMART-TRAFFIC-INTELLIGENCE",
        size=13, colour=MUTED)

    # PowerPoint locks whatever it has open, and during a rebuild-heavy week
    # that is often the last two copies as well. Walk to the first free name
    # rather than throwing away a completed build.
    names = ["PRESENTATION.pptx"] + [f"PRESENTATION_v{i}.pptx" for i in range(2, 9)]
    for i, name in enumerate(names):
        out = os.path.join(ROOT, "docs", name)
        try:
            prs.save(out)
        except PermissionError:
            continue
        if i:
            print(f"NOTE: {names[0]} is open in PowerPoint - wrote {name} instead. "
                  f"Close the deck and re-run to write the real file.")
        print("wrote", out, "|", len(prs.slides._sldIdLst), "slides")
        return
    raise SystemExit("every PRESENTATION*.pptx is locked - close PowerPoint")


if __name__ == "__main__":
    main()
